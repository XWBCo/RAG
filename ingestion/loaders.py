"""Custom document loaders for AlTi financial documents."""

import json
import re
from pathlib import Path
from typing import List, Optional

import pandas as pd
from llama_index.core import Document
from llama_index.readers.file import PyMuPDFReader


def preprocess_pdf_text(text: str) -> str:
    """
    Clean up PDF-extracted text for better chunking.

    Fixes common PDF extraction issues:
    - Merges lines broken mid-sentence
    - Adds sentence endings after headers
    - Normalizes whitespace
    - Preserves paragraph structure
    """
    if not text:
        return text

    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')

    # Remove page headers/footers (e.g., "AlTi Tiedemann Global    10")
    text = re.sub(r'AlTi Tiedemann Global\s+\d+\s*\n', '\n', text)

    # Add periods after common header patterns that lack punctuation
    # Pattern: Line that's all caps or title case, followed by newline
    header_patterns = [
        (r'\n(Purpose)\s*\n', r'\n\n\1.\n'),
        (r'\n(Overview)\s*\n', r'\n\n\1.\n'),
        (r'\n(How to Use)\s*\n', r'\n\n\1.\n'),
        (r'\n(Inputs & Outputs)\s*\n', r'\n\n\1.\n'),
        (r'\n(Example Use Case)\s*\n', r'\n\n\1.\n'),
        (r'\n(Best Uses)\s*\n', r'\n\n\1.\n'),
        (r'\n(Tips)\s*\n', r'\n\n\1.\n'),
        (r'\n(FAQs?)\s*\n', r'\n\n\1.\n'),
        (r'\n(Technical Glossary)\s*\n', r'\n\n\1.\n'),
    ]

    for pattern, replacement in header_patterns:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)

    # Add periods after "App – [Name]" headers
    text = re.sub(r'\n(App\s*[–-]\s*[A-Za-z\s]+)\s*\n', r'\n\n\1.\n\n', text)

    # Merge lines that were broken mid-sentence
    # If a line ends with a lowercase letter and next starts with lowercase, merge
    text = re.sub(r'([a-z,])\n([a-z])', r'\1 \2', text)

    # Preserve bullet points and list items - ensure they start new lines
    text = re.sub(r'\s*•\s*', '\n• ', text)
    text = re.sub(r'\n(\d+\.)\s+', r'\n\1 ', text)

    # Collapse multiple newlines to max 2 (paragraph break)
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Collapse multiple spaces
    text = re.sub(r'[ \t]+', ' ', text)

    # Clean up any remaining artifacts
    text = text.strip()

    return text

# PowerPoint support
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def load_powerpoint(file_path: Path, priority: str = "normal") -> List[Document]:
    """
    Load PowerPoint presentations - ONE DOCUMENT PER SLIDE.

    Creates separate documents for each slide to preserve slide-level
    context and improve retrieval granularity.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx required: pip install python-pptx")

    documents = []
    prs = Presentation(file_path)
    total_slides = len(prs.slides)

    # Priority score for retrieval boosting
    priority_scores = {"critical": 1.0, "high": 0.85, "normal": 0.5, "low": 0.3}

    for slide_idx, slide in enumerate(prs.slides):
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())

        # Skip empty slides
        if not slide_text_parts:
            continue

        slide_text = "\n".join(slide_text_parts)

        # Add context header for better retrieval
        full_text = f"""Document: {file_path.stem}
Slide {slide_idx + 1} of {total_slides}

{slide_text}
"""

        documents.append(Document(
            text=full_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "presentation_slide",
                "priority": priority,
                "priority_score": priority_scores.get(priority, 0.5),
                "slide_number": slide_idx + 1,
                "total_slides": total_slides,
                "is_priority_document": priority in ("critical", "high"),
            }
        ))

    return documents


def load_model_archetypes(file_path: Path) -> List[Document]:
    """
    Load Model Archetypes Excel file - PRIORITY DOCUMENT.

    This is a critical reference for investment model allocations.
    Creates rich, searchable documents for each fund and model.
    """
    documents = []
    xl = pd.ExcelFile(file_path, engine='openpyxl')

    # Priority metadata for all documents from this file
    base_metadata = {
        "source": str(file_path),
        "file_name": file_path.name,
        "priority": "critical",
        "priority_score": 1.0,
        "is_priority_document": True,
    }

    # 1. Process LIST sheet - Master fund database
    if "LIST" in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name="LIST")

        # Create individual fund documents for better retrieval
        for _, row in df.iterrows():
            if pd.isna(row.get("NAME")):
                continue

            fund_name = row.get("NAME", "Unknown")

            # Build comprehensive fund profile text
            fund_text = f"""Investment Fund Profile: {fund_name}

Risk Allocation: {row.get('RISK ALLOCATION', 'N/A')}
Asset Class: {row.get('ASSET CLASS', 'N/A')}
Structure: {row.get('STRUCTURE', 'N/A')}
Tax Treatment: {row.get('TAX ', 'N/A')}
Geography: {row.get('GEOGRAPHY', 'N/A')}
Return Target: {row.get('RETURN TARGET', 'N/A')}
Liquidity: {row.get('LIQUIDITY', 'N/A')}
Sizing: {row.get('SIZING', 'N/A')}
Role: {row.get('ROLE', 'N/A')}
Impact Theme: {row.get('IMPACT THEME', 'N/A')}

Description:
{row.get('DESCRIPTION', 'No description available.')}
"""
            documents.append(Document(
                text=fund_text,
                metadata={
                    **base_metadata,
                    "document_type": "fund_profile",
                    "sheet_name": "LIST",
                    "fund_name": fund_name,
                    "risk_allocation": str(row.get('RISK ALLOCATION', '')),
                    "asset_class": str(row.get('ASSET CLASS', '')),
                    "impact_theme": str(row.get('IMPACT THEME', '')),
                }
            ))

        # Also create a summary document listing all funds
        fund_summary = "Model Archetypes - Complete Fund Universe:\n\n"
        for _, row in df.iterrows():
            if pd.notna(row.get("NAME")):
                fund_summary += f"• {row['NAME']} ({row.get('ASSET CLASS', 'N/A')}) - {row.get('IMPACT THEME', 'N/A')}\n"

        documents.append(Document(
            text=fund_summary,
            metadata={
                **base_metadata,
                "document_type": "fund_universe_summary",
                "sheet_name": "LIST",
                "num_funds": len(df),
            }
        ))

    # 2. Process ALL MODELS sheets - hierarchy: Asset Class → Sub-Asset → Fund → Allocations
    risk_levels = ["CON", "BAL", "MG", "GRO", "LTG"]

    # Asset classes (appear in label column with totals, no fund name)
    asset_class_names = [
        "Stability Assets",
        "Diversified Assets",
        "Growth Assets - Public",
        "Growth Assets - Private",
        "Catalytic Debt",
        "Catalytic Equity",
    ]

    # Explicit model positions: (name, alias, label_col, fund_col, alloc_cols)
    # These positions are consistent across ALL MODELS and ALL MODELS (INT)
    MODEL_POSITIONS = [
        ("Integrated Best Ideas", "IBI", 1, 2, [3, 4, 5, 6, 7]),
        ("Impact 100%", "Impact", 10, 11, [12, 13, 14, 15, 16]),
        ("Climate Sustainability", "Climate", 19, 20, [21, 22, 23, 24, 25]),
        ("Inclusive Innovation", "Social", 28, 29, [30, 31, 32, 33, 34]),
        ("1099 Only", "Liquid", 37, 38, [39, 40, 41, 42, 43]),
    ]

    for sheet_name, region in [("ALL MODELS", "US"), ("ALL MODELS (INT)", "International")]:
        if sheet_name not in xl.sheet_names:
            continue

        df = pd.read_excel(xl, sheet_name=sheet_name, header=None)

        # Parse each model using explicit positions
        for model_name, model_alias, label_col, fund_col, alloc_cols in MODEL_POSITIONS:
            current_asset_class = ""
            current_sub_asset = ""

            # Row 5 is headers, row 6+ is data
            for row_idx in range(6, min(len(df), 55)):
                row = df.iloc[row_idx]

                # Get values safely
                label_val = row.iloc[label_col] if label_col < len(row) else None
                fund_val = row.iloc[fund_col] if fund_col < len(row) else None

                # Check for asset class row (has label, no fund name, has allocation totals)
                if pd.notna(label_val) and isinstance(label_val, str):
                    label = label_val.strip()
                    # Asset class rows have no fund name
                    if any(ac in label for ac in asset_class_names) and (pd.isna(fund_val) or not str(fund_val).strip()):
                        current_asset_class = label
                        current_sub_asset = ""
                        continue

                # Update sub-asset if label is present and not an asset class
                if pd.notna(label_val) and isinstance(label_val, str) and label_val.strip():
                    label = label_val.strip()
                    if not any(ac in label for ac in asset_class_names):
                        current_sub_asset = label

                # Get fund name - skip if missing
                if pd.isna(fund_val) or not isinstance(fund_val, str):
                    continue
                fund_name = fund_val.strip()
                if not fund_name or fund_name.lower() in ["nan", "", "-", "tbd"]:
                    continue

                # Get allocations [CON, BAL, MG, GRO, LTG]
                allocs = []
                for ac in alloc_cols:
                    if ac < len(row):
                        val = row.iloc[ac]
                        if pd.notna(val) and isinstance(val, (int, float)):
                            allocs.append(float(val))
                        else:
                            allocs.append(0.0)
                    else:
                        allocs.append(0.0)

                # Skip if all zeros (not allocated in this model)
                if all(a == 0 for a in allocs):
                    continue

                # Format as percentages
                alloc_pcts = [f"{a*100:.1f}" for a in allocs]
                alloc_array_str = f"[{', '.join(alloc_pcts)}]"

                # Create document with rich searchable text
                doc_text = f"""Fund Allocation: {fund_name}

Model: {model_name} ({model_alias})
Region: {region}
Asset Class: {current_asset_class}
Sub-Asset: {current_sub_asset}
Fund: {fund_name}

Allocations % [CON, BAL, MG, GRO, LTG]: {alloc_array_str}

"""
                # Searchable statements for each risk level
                for i, risk in enumerate(risk_levels):
                    pct = allocs[i] * 100
                    if pct > 0:
                        doc_text += f"In {model_name} at {risk} risk, {fund_name} is allocated {pct:.1f}%.\n"

                # Additional searchable context
                if current_sub_asset:
                    doc_text += f"\n{fund_name} serves as the {current_sub_asset.lower()} allocation in {model_name}'s {current_asset_class.lower()}.\n"

                documents.append(Document(
                    text=doc_text,
                    metadata={
                        **base_metadata,
                        "document_type": "fund_model_allocation",
                        "model_name": model_name,
                        "model_alias": model_alias,
                        "model_region": region,
                        "asset_class": current_asset_class,
                        "sub_asset": current_sub_asset,
                        "fund_name": fund_name,
                        "alloc_con": allocs[0],
                        "alloc_bal": allocs[1],
                        "alloc_mg": allocs[2],
                        "alloc_gro": allocs[3],
                        "alloc_ltg": allocs[4],
                        "is_liquid_model": model_name == "1099 Only",
                    }
                ))

    # Overview document with all 5 models
    documents.append(Document(
        text="""Portfolio Model Archetypes Overview

Models (Investment Menus):
1. Integrated Best Ideas (IBI) - Returns-first with ESG screening
2. Impact 100% - Comprehensive measurable impact across all themes
3. Climate Sustainability (Climate) - Environment-first investment approach
4. Inclusive Innovation (Social) - Social equity and inclusive economic growth
5. 1099 Only (Liquid) - K1-free liquid variant for tax-sensitive clients

Risk Levels: CON (Conservative), BAL (Balanced), MG (Moderate Growth), GRO (Growth), LTG (Long-Term Growth)

Regions: US (domestic), International (non-US)

Asset Class Hierarchy:
- Stability Assets: Cash Equivalent, IG Fixed Income
- Diversified Assets: Real Estate, Infrastructure, Private Credit, Sust. Forestry
- Growth Assets - Public: US Passive, Global Active, Non US, Thematic, Global EM
- Growth Assets - Private: Vintages (GIO, PEP), Direct investments
- Catalytic Debt: CDFIs (California Farmlink, Craft3, Oweesta, Seed Commons, Akiptan)
- Catalytic Equity: Impact-focused private investments

Client model selection: {Archetype} + {Region} + {Tax Preference} + {Risk Level}
Allocations vary by model AND risk level. Use format: [CON, BAL, MG, GRO, LTG]
""",
        metadata={
            **base_metadata,
            "document_type": "model_overview",
        }
    ))

    return documents


def load_portfolio_csv(file_path: Path, portfolio_name: Optional[str] = None) -> List[Document]:
    """
    Load portfolio holdings from CSV files.

    Converts structured portfolio data into documents with rich metadata
    for semantic search on fund names, tickers, and allocations.

    KEY: Groups by Portfolio/Benchmark Name to preserve model context.
    """
    documents = []

    try:
        df = pd.read_csv(file_path, encoding="utf-8")
    except UnicodeDecodeError:
        df = pd.read_csv(file_path, encoding="latin-1")

    # Check if file has Portfolio Name or Benchmark Name column (grouped data)
    name_col = None
    for col in ["Portfolio Name", "Benchmark Name", "Model Name", "Model"]:
        if col in df.columns:
            name_col = col
            break

    if name_col and "Weight" in df.columns:
        # GROUP BY PORTFOLIO/MODEL - create one document per model with complete holdings
        for model_name, group in df.groupby(name_col):
            # Sort by weight descending
            group_sorted = group.sort_values("Weight", ascending=False)

            # Build complete holdings list
            holdings_lines = []
            total_weight = 0.0
            for _, row in group_sorted.iterrows():
                fund = row.get("Tier4", row.get("Name", "Unknown"))
                weight = row.get("Weight", 0)
                total_weight += weight
                holdings_lines.append(f"- {fund}: {weight:.2%}")

            holdings_text = "\n".join(holdings_lines)

            # Create comprehensive document for this model
            doc_text = f"""Portfolio Model: {model_name}

COMPLETE ALLOCATION BREAKDOWN:
Total Holdings: {len(group)}
Total Weight: {total_weight:.2%}

All Holdings (sorted by weight):
{holdings_text}

This is the complete allocation for the {model_name} portfolio model.
"""
            documents.append(Document(
                text=doc_text,
                metadata={
                    "source": str(file_path),
                    "file_name": file_path.name,
                    "document_type": "portfolio_model_complete",
                    "portfolio_name": model_name,
                    "model_name": model_name,
                    "num_holdings": len(group),
                    "total_weight": total_weight,
                }
            ))

            # Also create a summary for quick retrieval
            top_5 = group_sorted.head(5)
            top_5_text = "\n".join([
                f"- {row.get('Tier4', 'Unknown')}: {row['Weight']:.2%}"
                for _, row in top_5.iterrows()
            ])

            summary_text = f"""Portfolio: {model_name}
Top 5 Holdings:
{top_5_text}

Total holdings in {model_name}: {len(group)}
"""
            documents.append(Document(
                text=summary_text,
                metadata={
                    "source": str(file_path),
                    "file_name": file_path.name,
                    "document_type": "portfolio_summary",
                    "portfolio_name": model_name,
                    "model_name": model_name,
                    "num_holdings": len(group),
                }
            ))

        return documents

    # Fallback: No grouping column - use filename as portfolio name
    if portfolio_name is None:
        portfolio_name = file_path.stem.replace("_", " ").title()

    # Create a summary document for the entire portfolio
    if "Weight" in df.columns:
        top_holdings = df.nlargest(10, "Weight") if len(df) > 10 else df
        holdings_text = "\n".join([
            f"- {row.get('Tier4', row.get('Name', 'Unknown'))}: {row['Weight']:.2%}"
            for _, row in top_holdings.iterrows()
        ])

        summary = f"""Portfolio: {portfolio_name}
Total Holdings: {len(df)}
Top Holdings:
{holdings_text}
"""
        documents.append(Document(
            text=summary,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "portfolio_summary",
                "portfolio_name": portfolio_name,
                "num_holdings": len(df),
            }
        ))

    # Create documents for each holding (chunked by groups of 20)
    chunk_size = 20
    for i in range(0, len(df), chunk_size):
        chunk = df.iloc[i:i + chunk_size]
        holdings_text = "\n".join([
            f"{row.get('Tier4', row.get('Name', 'Unknown'))}: "
            f"{row.get('Weight', 0):.2%} allocation"
            for _, row in chunk.iterrows()
        ])

        documents.append(Document(
            text=f"Portfolio: {portfolio_name}\nHoldings (Part {i // chunk_size + 1}):\n{holdings_text}",
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "portfolio_holdings",
                "portfolio_name": portfolio_name,
                "chunk_index": i // chunk_size,
            }
        ))

    return documents


def load_cma_excel(file_path: Path) -> List[Document]:
    """
    Load Capital Market Assumptions from Excel files.

    Handles multi-sheet Excel files with returns, correlations,
    and time series data.
    """
    documents = []
    xl = pd.ExcelFile(file_path)

    for sheet_name in xl.sheet_names:
        try:
            df = pd.read_excel(xl, sheet_name=sheet_name)

            # Skip empty sheets
            if df.empty:
                continue

            # Create document based on sheet type
            if sheet_name.upper() in ["RR", "RETURNS_RISK", "ASSUMPTIONS"]:
                # Returns and Risk sheet - asset class assumptions
                text_parts = [f"Capital Market Assumptions - {sheet_name}:\n"]

                for _, row in df.iterrows():
                    asset_class = row.iloc[0] if len(row) > 0 else "Unknown"
                    values = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    text_parts.append(f"- {asset_class}: {values}")

                documents.append(Document(
                    text="\n".join(text_parts),
                    metadata={
                        "source": str(file_path),
                        "file_name": file_path.name,
                        "sheet_name": sheet_name,
                        "document_type": "cma_assumptions",
                    }
                ))

            elif sheet_name.upper() in ["CORR", "CORRELATION", "CORRELATIONS"]:
                # Correlation matrix - summarize key relationships
                text = f"Correlation Matrix ({sheet_name}):\n"
                text += f"Asset classes covered: {', '.join(df.columns[:10].tolist())}...\n"
                text += f"Matrix dimensions: {df.shape[0]} x {df.shape[1]}"

                documents.append(Document(
                    text=text,
                    metadata={
                        "source": str(file_path),
                        "file_name": file_path.name,
                        "sheet_name": sheet_name,
                        "document_type": "cma_correlation",
                    }
                ))

            else:
                # Generic sheet - convert to text
                text = f"CMA Data - {sheet_name}:\n"
                text += df.head(20).to_string()

                documents.append(Document(
                    text=text,
                    metadata={
                        "source": str(file_path),
                        "file_name": file_path.name,
                        "sheet_name": sheet_name,
                        "document_type": "cma_data",
                    }
                ))

        except Exception as e:
            print(f"Warning: Could not process sheet {sheet_name}: {e}")
            continue

    return documents


def load_returns_csv(file_path: Path) -> List[Document]:
    """
    Load historical returns time series from CSV.

    Creates summary documents with key statistics rather than
    vectorizing raw time series data.
    """
    documents = []

    try:
        df = pd.read_csv(file_path, index_col=0, parse_dates=True)
    except Exception:
        df = pd.read_csv(file_path)

    # Create summary statistics document
    if df.select_dtypes(include=["float64", "int64"]).shape[1] > 0:
        numeric_df = df.select_dtypes(include=["float64", "int64"])

        stats_text = f"Historical Returns Summary ({file_path.stem}):\n"
        stats_text += f"Date Range: {df.index[0]} to {df.index[-1]}\n" if hasattr(df.index[0], "strftime") else ""
        stats_text += f"Asset Classes: {', '.join(numeric_df.columns.tolist())}\n\n"

        for col in numeric_df.columns[:15]:  # Limit to 15 assets
            mean_return = numeric_df[col].mean()
            volatility = numeric_df[col].std()
            stats_text += f"- {col}: Mean {mean_return:.2%}, Vol {volatility:.2%}\n"

        documents.append(Document(
            text=stats_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "returns_summary",
                "num_periods": len(df),
                "num_assets": len(numeric_df.columns),
            }
        ))

    return documents


def load_pdf_documents(file_path: Path) -> List[Document]:
    """
    Load PDF documents (research reports, client documents).

    Uses PyMuPDF for robust PDF parsing with page-level metadata.
    Preprocesses text to improve chunking quality.
    """
    reader = PyMuPDFReader()
    raw_documents = reader.load(file_path=file_path)

    documents = []
    for i, doc in enumerate(raw_documents):
        # Preprocess text for better chunking
        cleaned_text = preprocess_pdf_text(doc.text)

        # Skip empty pages
        if not cleaned_text or len(cleaned_text.strip()) < 50:
            continue

        documents.append(Document(
            text=cleaned_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "pdf_document",
                "page_number": i + 1,
                "priority": "normal",
            }
        ))

    return documents


def load_markdown_documents(file_path: Path) -> List[Document]:
    """
    Load markdown documents (FAQ, documentation, guides).

    Splits by section headers (##) to create logical chunks.
    """
    documents = []

    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    # Get document title from first line
    lines = content.split("\n")
    title = lines[0].lstrip("#").strip() if lines else file_path.stem

    # Split by ## headers to create logical sections
    sections = []
    current_section = {"header": title, "content": []}

    for line in lines[1:]:
        if line.startswith("## "):
            # Save previous section
            if current_section["content"]:
                sections.append(current_section)
            current_section = {"header": line.lstrip("#").strip(), "content": []}
        else:
            current_section["content"].append(line)

    # Don't forget the last section
    if current_section["content"]:
        sections.append(current_section)

    # Create documents from sections
    for i, section in enumerate(sections):
        section_text = "\n".join(section["content"]).strip()
        if not section_text or len(section_text) < 50:
            continue

        # Include header in text for better retrieval
        full_text = f"{section['header']}\n\n{section_text}"

        documents.append(Document(
            text=full_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "faq_section",
                "document_title": title,
                "section_header": section["header"],
                "section_index": i,
                "priority": "high",  # FAQ content is high priority
                "priority_score": 0.85,
            }
        ))

    return documents


def load_qualtrics_json(file_path: Path) -> List[Document]:
    """
    Load Qualtrics survey responses from JSON files.

    Extracts client information and survey responses for
    investment preference analysis.
    """
    documents = []

    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Handle different JSON structures
    if isinstance(data, list):
        responses = data
    elif isinstance(data, dict) and "responses" in data:
        responses = data["responses"]
    else:
        responses = [data]

    for i, response in enumerate(responses):
        # Extract client info if available
        client_name = response.get("client_name", response.get("ClientName", f"Client {i + 1}"))

        # Build response text
        text_parts = [f"Survey Response - {client_name}:\n"]

        for key, value in response.items():
            if key.startswith("Q") or key.lower() in ["risk_tolerance", "investment_horizon", "goals"]:
                text_parts.append(f"- {key}: {value}")

        documents.append(Document(
            text="\n".join(text_parts),
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "survey_response",
                "client_name": client_name,
                "response_index": i,
            }
        ))

    return documents


def load_fund_holdings_csv(file_path: Path) -> List[Document]:
    """
    Load fund holdings files (SMA, INTL, LP holdings).

    These are larger files with individual security positions.
    """
    documents = []

    try:
        df = pd.read_csv(file_path, encoding="utf-8")
    except UnicodeDecodeError:
        df = pd.read_csv(file_path, encoding="latin-1")

    fund_type = file_path.stem.split("_")[0].upper()

    # Summary document
    summary = f"""Fund Holdings: {fund_type}
Total Positions: {len(df)}
File: {file_path.name}
"""

    if "Weight" in df.columns:
        total_weight = df["Weight"].sum()
        summary += f"Total Weight: {total_weight:.2%}\n"

    documents.append(Document(
        text=summary,
        metadata={
            "source": str(file_path),
            "file_name": file_path.name,
            "document_type": "fund_holdings_summary",
            "fund_type": fund_type,
            "num_positions": len(df),
        }
    ))

    # Chunk holdings for searchability
    chunk_size = 50
    for i in range(0, len(df), chunk_size):
        chunk = df.iloc[i:i + chunk_size]

        holdings_text = []
        for _, row in chunk.iterrows():
            name = row.get("Name", row.get("Security", row.iloc[0]))
            weight = row.get("Weight", "N/A")
            if isinstance(weight, float):
                weight = f"{weight:.4%}"
            holdings_text.append(f"- {name}: {weight}")

        documents.append(Document(
            text=f"{fund_type} Holdings (Positions {i + 1}-{i + len(chunk)}):\n" + "\n".join(holdings_text),
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "document_type": "fund_holdings_detail",
                "fund_type": fund_type,
                "chunk_index": i // chunk_size,
            }
        ))

    return documents
